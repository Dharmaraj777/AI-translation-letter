import io
from typing import List, Dict, Optional

from pptx import Presentation  # python-pptx
from pptx.enum.shapes import MSO_SHAPE_TYPE

from spanish_translator_oai_client import OaiClient
from spanish_translator_logger import logger


class PptxProcessor:
    def __init__(self, oai_client: OaiClient):
        self.oai_client = oai_client

    """
    PPTX translator.

    Key points:
    - We translate at RUN level inside paragraphs of:
        - shape.text_frame.paragraphs[].runs
        - table.cell.text_frame.paragraphs[].runs
        - shapes nested inside GROUP shapes (recursive)
    - This preserves text formatting (font, size, color, etc.) and does not
      remove shapes or images, because we never reassign paragraph.text or cell.text.
    - Speaker notes are intentionally NOT translated.
    """

    # ------------------------------------------------------------------
    # Internal: collect text from a single shape (recursively)
    # ------------------------------------------------------------------
    def _collect_from_shape(
        self,
        shape,
        segments: List[Dict[str, str]],
        s_idx: int,
        shape_path: str,
    ) -> None:
        """
        Recursively collect run-level text segments from a shape and any nested shapes
        (especially group shapes).
        """
        # 1) Normal shapes with text
        if getattr(shape, "has_text_frame", False):
            text_frame = shape.text_frame
            if text_frame is not None:
                for p_idx, paragraph in enumerate(text_frame.paragraphs):
                    for r_idx, run in enumerate(paragraph.runs):
                        text = run.text
                        if text and text.strip():
                            seg_id = f"s-{s_idx}-{shape_path}-p-{p_idx}-r-{r_idx}"
                            segments.append({"id": seg_id, "text": text})

        # 2) Tables inside shapes
        if getattr(shape, "has_table", False):
            table = shape.table
            if table is not None:
                for row_idx, row in enumerate(table.rows):
                    for col_idx, cell in enumerate(row.cells):
                        text_frame = cell.text_frame
                        if text_frame is None:
                            continue
                        for p_idx, paragraph in enumerate(text_frame.paragraphs):
                            for r_idx, run in enumerate(paragraph.runs):
                                text = run.text
                                if text and text.strip():
                                    seg_id = (
                                        f"s-{s_idx}-{shape_path}-tbl-row-"
                                        f"{row_idx}-col-{col_idx}-p-{p_idx}-r-{r_idx}"
                                    )
                                    segments.append({"id": seg_id, "text": text})

        # 3) Group shapes → recurse into children
        # (They have shape_type == GROUP and a .shapes collection.)
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for g_idx, subshape in enumerate(shape.shapes):
                    sub_path = f"{shape_path}-g-{g_idx}"
                    self._collect_from_shape(subshape, segments, s_idx, sub_path)
        except Exception:
            # Some shapes may not expose .shape_type or .shapes cleanly
            pass

    # ------------------------------------------------------------------
    # Segment collection (for entire presentation)
    # ------------------------------------------------------------------
    def _collect_segments(self, pres: Presentation) -> List[Dict[str, str]]:
        """
        Collect text segments at RUN level from all slides, shapes, table cells,
        and nested group shapes.

        Each segment has:
          { "id": <segment_id>, "text": <original_text> }
        """
        segments: List[Dict[str, str]] = []

        for s_idx, slide in enumerate(pres.slides):
            for sh_idx, shape in enumerate(slide.shapes):
                shape_path = f"sh-{sh_idx}"
                self._collect_from_shape(shape, segments, s_idx, shape_path)

        logger.info(f"[pptx] Collected {len(segments)} text segments from PPTX.")
        return segments

    # ------------------------------------------------------------------
    # Internal: apply translations to a single shape (recursively)
    # ------------------------------------------------------------------
    def _apply_to_shape(
        self,
        shape,
        id_to_translation: Dict[str, str],
        s_idx: int,
        shape_path: str,
    ) -> None:
        """
        Recursively walk a shape (and nested shapes) and update run.text
        where we have translations.
        """

        # 1) Normal shape text
        if getattr(shape, "has_text_frame", False):
            text_frame = shape.text_frame
            if text_frame is not None:
                for p_idx, paragraph in enumerate(text_frame.paragraphs):
                    for r_idx, run in enumerate(paragraph.runs):
                        text = run.text
                        if not text or not text.strip():
                            continue
                        seg_id = f"s-{s_idx}-{shape_path}-p-{p_idx}-r-{r_idx}"
                        if seg_id in id_to_translation:
                            run.text = id_to_translation[seg_id]

        # 2) Table text
        if getattr(shape, "has_table", False):
            table = shape.table
            if table is not None:
                for row_idx, row in enumerate(table.rows):
                    for col_idx, cell in enumerate(row.cells):
                        text_frame = cell.text_frame
                        if text_frame is None:
                            continue
                        for p_idx, paragraph in enumerate(text_frame.paragraphs):
                            for r_idx, run in enumerate(paragraph.runs):
                                text = run.text
                                if not text or not text.strip():
                                    continue
                                seg_id = (
                                    f"s-{s_idx}-{shape_path}-tbl-row-"
                                    f"{row_idx}-col-{col_idx}-p-{p_idx}-r-{r_idx}"
                                )
                                if seg_id in id_to_translation:
                                    run.text = id_to_translation[seg_id]

        # 3) Group shapes → recurse into children
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for g_idx, subshape in enumerate(shape.shapes):
                    sub_path = f"{shape_path}-g-{g_idx}"
                    self._apply_to_shape(subshape, id_to_translation, s_idx, sub_path)
        except Exception:
            pass

    # ------------------------------------------------------------------
    # Apply translations for the entire presentation
    # ------------------------------------------------------------------
    def _apply_translations(
        self,
        pres: Presentation,
        id_to_translation: Dict[str, str],
    ) -> None:
        """
        Walk the PPT structure again and update run.text where we have translations.
        """
        for s_idx, slide in enumerate(pres.slides):
            for sh_idx, shape in enumerate(slide.shapes):
                shape_path = f"sh-{sh_idx}"
                self._apply_to_shape(shape, id_to_translation, s_idx, shape_path)

    # ------------------------------------------------------------------
    # Public entrypoint
    # ------------------------------------------------------------------
    def translate_document(
        self,
        filename: str,
        content_bytes: bytes,
        target_language: Optional[str] = None,
        target_dialect: Optional[str] = None,
    ) -> bytes:
        logger.info(f"Translating PPTX document: {filename}")
        pres_stream = io.BytesIO(content_bytes)
        pres = Presentation(pres_stream)

        # 1) Collect run-level segments (including nested group shapes)
        segments = self._collect_segments(pres)
        if not segments:
            logger.info("No text segments found in PPTX, returning original document.")
            return content_bytes

        # 2) Call Azure OpenAI to translate
        id_to_translation = self.oai_client.translate_segments(
            segments,
            target_language=target_language,
            target_dialect=target_dialect,
        )

        # 3) Apply translations
        self._apply_translations(pres, id_to_translation)

        # 4) Save to bytes
        out_stream = io.BytesIO()
        pres.save(out_stream)
        out_stream.seek(0)
        return out_stream.read()
