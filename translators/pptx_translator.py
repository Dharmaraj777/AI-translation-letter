import io
from typing import List, Dict, Optional

from pptx import Presentation  # python-pptx

from .base_translator import BaseTranslator
from ai_translation_logger import logger
from ai_translation_utils import UtilityFunctions


class PptxTranslator(BaseTranslator):
    """
    PPTX translator.

    Key points:
    - We translate at RUN level inside paragraphs of:
        - shape.text_frame.paragraphs[].runs
        - table.cell.text_frame.paragraphs[].runs
    - This preserves text formatting (font, size, color, etc.) and does not
      remove shapes or images, because we never reassign paragraph.text or cell.text.
    """

    def can_handle(self, filename: str) -> bool:
        return UtilityFunctions.get_extension(filename) == ".pptx"

    # ------------------------------------------------------------------
    # Segment collection
    # ------------------------------------------------------------------
    def _collect_segments(self, pres: Presentation) -> List[Dict[str, str]]:
        """
        Collect text segments at RUN level from all slides, shapes, and table cells.
        Each segment has:
          { "id": <segment_id>, "text": <original_text> }
        """
        segments: List[Dict[str, str]] = []

        for s_idx, slide in enumerate(pres.slides):
            for sh_idx, shape in enumerate(slide.shapes):
                # 1) Normal shapes with text
                if shape.has_text_frame:
                    text_frame = shape.text_frame
                    for p_idx, paragraph in enumerate(text_frame.paragraphs):
                        for r_idx, run in enumerate(paragraph.runs):
                            text = run.text
                            if text and text.strip():
                                seg_id = (
                                    f"s-{s_idx}-sh-{sh_idx}-p-{p_idx}-r-{r_idx}"
                                )
                                segments.append({"id": seg_id, "text": text})

                # 2) Tables inside shapes
                if shape.has_table:
                    table = shape.table
                    for row_idx, row in enumerate(table.rows):
                        for col_idx, cell in enumerate(row.cells):
                            # cell.text_frame holds the paragraphs
                            text_frame = cell.text_frame
                            for p_idx, paragraph in enumerate(text_frame.paragraphs):
                                for r_idx, run in enumerate(paragraph.runs):
                                    text = run.text
                                    if text and text.strip():
                                        seg_id = (
                                            f"s-{s_idx}-sh-{sh_idx}-tbl-row-"
                                            f"{row_idx}-col-{col_idx}-p-{p_idx}-r-{r_idx}"
                                        )
                                        segments.append({"id": seg_id, "text": text})

        return segments

    # ------------------------------------------------------------------
    # Apply translations back to the presentation
    # ------------------------------------------------------------------
    def _apply_translations(
        self,
        pres: Presentation,
        id_to_translation: Dict[str, str],
    ) -> None:
        """
        Walk the PPT structure again and update run.text where we have translations.
        """

        for s_idx, slide in enumerate(pres.slides):
            for sh_idx, shape in enumerate(slide.shapes):

                # 1) Normal shape text
                if shape.has_text_frame:
                    text_frame = shape.text_frame
                    for p_idx, paragraph in enumerate(text_frame.paragraphs):
                        for r_idx, run in enumerate(paragraph.runs):
                            text = run.text
                            if not text or not text.strip():
                                continue
                            seg_id = (
                                f"s-{s_idx}-sh-{sh_idx}-p-{p_idx}-r-{r_idx}"
                            )
                            if seg_id in id_to_translation:
                                run.text = id_to_translation[seg_id]

                # 2) Table text
                if shape.has_table:
                    table = shape.table
                    for row_idx, row in enumerate(table.rows):
                        for col_idx, cell in enumerate(row.cells):
                            text_frame = cell.text_frame
                            for p_idx, paragraph in enumerate(text_frame.paragraphs):
                                for r_idx, run in enumerate(paragraph.runs):
                                    text = run.text
                                    if not text or not text.strip():
                                        continue
                                    seg_id = (
                                        f"s-{s_idx}-sh-{sh_idx}-tbl-row-"
                                        f"{row_idx}-col-{col_idx}-p-{p_idx}-r-{r_idx}"
                                    )
                                    if seg_id in id_to_translation:
                                        run.text = id_to_translation[seg_id]

    # ------------------------------------------------------------------
    # Public entrypoint
    # ------------------------------------------------------------------
    def translate_document(
        self,
        filename: str,
        content_bytes: bytes,
        target_language: Optional[str] = None,
        target_dialect: Optional[str] = None,
    ) -> bytes:
        logger.info(f"Translating PPTX document: {filename}")
        pres_stream = io.BytesIO(content_bytes)
        pres = Presentation(pres_stream)

        # 1) Collect run-level segments
        segments = self._collect_segments(pres)
        if not segments:
            logger.info("No text segments found in PPTX, returning original document.")
            return content_bytes

        # 2) Call Azure OpenAI to translate
        id_to_translation = self.oai_client.translate_segments(
            segments,
            target_language=target_language,
            target_dialect=target_dialect,
        )

        # 3) Apply translations
        self._apply_translations(pres, id_to_translation)

        # 4) Save to bytes
        out_stream = io.BytesIO()
        pres.save(out_stream)
        out_stream.seek(0)
        return out_stream.read()
